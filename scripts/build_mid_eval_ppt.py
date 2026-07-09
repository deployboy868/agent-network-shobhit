"""Generate mid-internship evaluation PowerPoint."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "presentation" / "Mid_Intern_Evaluation_Agent_Network.pptx"

# Sprinklr-adjacent professional palette
NAVY = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT = RGBColor(0x03, 0x84, 0xC6)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.15))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(9), Inches(0.55))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(9), Inches(0.35))
        stf = sub.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)


def _add_bullets(slide, items: list[str], left=0.55, top=1.45, width=8.9, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = SLATE
        p.space_after = Pt(10)


def _add_table_slide(slide, headers: list[str], rows: list[list[str]]):
    rows_n = len(rows) + 1
    cols_n = len(headers)
    table_shape = slide.shapes.add_table(rows_n, cols_n, Inches(0.45), Inches(1.5), Inches(9.1), Inches(0.4 + 0.38 * rows_n))
    table = table_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = SLATE


def build() -> Path:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)
    title = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(8.6), Inches(1.2))
    tf = title.text_frame
    tf.text = "Agent Social Network"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(1.5))
    stf = sub.text_frame
    stf.text = "Mid-Internship Evaluation\nShobhit Raj · IIT Patna · Engineering Intern\nMentors: Mayank Kumar, Prabhpreet Singh · Pillar: Team Productivity"
    for i, para in enumerate(stf.paragraphs):
        para.font.size = Pt(16 if i == 0 else 14)
        para.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        para.space_after = Pt(6)

    accent = slide.shapes.add_shape(1, Inches(0.7), Inches(5.0), Inches(1.8), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # --- Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "Today’s Agenda", "Aligned with mid-evaluation review")
    _add_bullets(
        slide,
        [
            "1. Problem statement",
            "2. Why this is a real problem",
            "3. Who our users are",
            "4. Current progress & live demos",
            "5. End vision for the product",
            "6. Plan for upcoming weeks",
            "7. Blockers & how we’re addressing them",
        ],
        size=20,
    )

    # --- 1 Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "1. Problem Statement")
    _add_bullets(
        slide,
        [
            "Build an internal Agent Social Network at Sprinklr.",
            "Each node = a digital twin of an employee.",
            "Twins use MCP tools: Jira, GitLab, Teams, Workday.",
            "Twins coordinate with each other to assign work and follow up.",
            "Goal: remove routine context-switching between tools and people.",
        ],
    )

    # --- 2 Why ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "2. Why This Matters")
    _add_bullets(
        slide,
        [
            "Managers assign in Jira — engineers may not see it immediately.",
            "Status checks require manual lookups across Jira, GitLab, Teams.",
            "MR context lives in GitLab; ticket context lives in Jira.",
            "Cross-team coordination burns time on copying context, not building.",
            "Agents + agent-to-agent messaging automate the routine coordination layer.",
        ],
    )

    # --- 3 Users ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "3. User Base")
    _add_bullets(
        slide,
        [
            "Managers / leads — assign work, review team progress via their twin.",
            "Engineers — receive tasks, update Jira, complete work via their twin.",
            "Platform / intern team — build twins, MCP tools, orchestration, safety.",
            "Today: demo twins (Manager, Assignee, Observer) in code.",
            "End state: every employee from Workday roster gets a twin.",
        ],
    )

    # --- 4 Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "4. Current Progress — Architecture")
    arch = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(8.9), Inches(5.3))
    tf = arch.text_frame
    tf.text = (
        "User / Demo  →  LangGraph (delegate → complete → track)\n"
        "                    ↓\n"
        "           Digital Twin Agents  ↔  Agent Message Bus\n"
        "                    ↓\n"
        "              MCP Tool Layer\n"
        "         Jira (live) · GitLab (read) · Teams · Workday\n"
        "                    ↓\n"
        "           MCP Server in Cursor (10 tools, live)\n\n"
        "Next layer (planned): LLM reasoning on top of MCP + LangGraph"
    )
    for p in tf.paragraphs:
        p.font.name = "Menlo"
        p.font.size = Pt(15)
        p.font.color.rgb = SLATE

    # --- 4b Progress table ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "4. Current Progress — Deliverables")
    _add_table_slide(
        slide,
        ["Deliverable", "Status", "Notes"],
        [
            ["Digital twins + skills", "Done", "Sample roster; Workday next"],
            ["Jira MCP", "Live", "Create, assign, close, review"],
            ["GitLab MCP", "Partial", "Read-only + Jira comment link"],
            ["Teams / Workday", "Mock", "Awaiting mentor access"],
            ["Agent-to-agent bus", "Done", "In-process protocol MVP"],
            ["Assign + track demo", "Live", "Safe-prefix + safe mode"],
            ["LangGraph", "Done", "delegate → complete → track"],
            ["MCP server (Cursor)", "Live", "Green — 10 tools"],
            ["LLM autonomy", "Planned", "MCP ready as tool layer"],
        ],
    )

    # --- 4c Demos ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "4. Live Demos You Can Show")
    _add_bullets(
        slide,
        [
            "assign_and_track_graph — full live Jira flow via LangGraph",
            "review_tasks — read-only manager progress view",
            "MCP in Cursor — agent_network_status, jira_list_tickets",
            "verify_gitlab — read-only MR list (VPN)",
            "Safety: [Agent-Network-TEST] prefix · demo assigns to intern only",
        ],
        size=17,
    )

    # --- 5 Vision ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "5. End Vision")
    _add_bullets(
        slide,
        [
            "Every employee has a digital twin with skills from their role.",
            "Manager: “Assign the handbook fix to the assignee” (natural language).",
            "LLM twin picks tools, creates Jira ticket, messages assignee’s twin.",
            "Assignee twin works ticket; GitLab MR linked; Teams notification sent.",
            "Reporter twin autonomously follows up until done — no manual chasing.",
            "MCP is the shared tool layer for Cursor, Copilot, and internal agents.",
        ],
    )

    # --- 6 Plan ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "6. Plan — Upcoming Weeks")
    _add_table_slide(
        slide,
        ["Week", "Focus", "Outcome"],
        [
            ["Week 3", "Teams integration", "Live notify assignee (test channel)"],
            ["Week 3", "Workday CSV loader", "Real employee roster → twins"],
            ["Week 3–4", "GitLab in LangGraph", "Link MR → Jira in main flow"],
            ["Week 4", "LLM + tool calling", "Natural language → MCP tools"],
            ["Week 4", "Final demo", "End-to-end story for evaluation"],
        ],
    )

    # --- 7 Blockers ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_header_bar(slide, "7. Blockers & Mitigations")
    _add_table_slide(
        slide,
        ["Blocker", "Impact", "Mitigation"],
        [
            ["Teams API / channel", "No live pings", "Mock works; requested access"],
            ["Workday roster", "Sample twins only", "CSV loader ready"],
            ["GitLab UI blocked", "Can’t browse MRs", "API read works with VPN"],
            ["No open MRs in main repo", "Can’t demo MR link", "Code ready"],
            ["LLM API choice", "No NL autonomy yet", "Ask Sprinklr-approved API"],
        ],
    )

    # --- Close ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(2))
    tf = box.text_frame
    tf.text = "Thank you\nQuestions & feedback welcome"
    for i, p in enumerate(tf.paragraphs):
        p.font.size = Pt(36 if i == 0 else 20)
        p.font.bold = i == 0
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
